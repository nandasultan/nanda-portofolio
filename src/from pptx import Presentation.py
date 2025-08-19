from pptx import Presentation
from pptx.util import Inches

# Buat presentasi baru
prs = Presentation()

# Layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Fungsi bantu tambah slide
def add_slide(title, content):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text = content

# Slide 1 – Judul
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Belajar Angklung Singkat & Praktik"
slide.placeholders[1].text = "Nama Guru – Kelas – Tanggal"

# Slide 2 – Apa itu Angklung
add_slide("Apa itu Angklung",
          "- Alat musik tradisional dari bambu\n- Dimainkan dengan cara digoyangkan\n- Menghasilkan satu nada per angklung")

# Slide 3 – Cara Memainkan Angklung
add_slide("Cara Memainkan Angklung",
          "- Pegang bagian atas\n- Goyangkan ke satu arah dengan stabil\n- Mainkan saat aba-aba atau isyarat")

# Slide 4 – Apa itu Melodi & Contoh Alat Musik Melodi
add_slide("Apa itu Melodi & Contoh Alat Musik Melodi",
          "Melodi: susunan nada yang membentuk lagu\n\nContoh alat musik melodi:\n1. Angklung\n2. Pianika\n3. Gitar\n4. Recorder")

# Slide 5 – Do Re Mi & Angka
add_slide("Do Re Mi & Angka",
          "Do = 1\nRe = 2\nMi = 3\nFa = 4\nSol = 5\nLa = 6\nSi = 7\n\nAngklung biasa pakai angka untuk nada")

# Slide 6 – Akor C
add_slide("Akor C",
          "C = Do = 1\nAkor C = 1 – 3 – 5\n(Do – Mi – Sol)")

# Slide 7 – Praktik
add_slide("Praktik",
          "🎶 Saatnya bermain angklung!\n- Ikuti aba-aba\n- Coba mainkan Do – Mi – Sol\n- Lagu sederhana: Balonku atau Twinkle-Twinkle")

# Simpan file
prs.save("Belajar_Angklung_Singkat.pptx")
